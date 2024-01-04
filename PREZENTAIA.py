#-----------------04.12.23
#-----------------rabotaet cmeste s shablonom i tmn rm 14
import os
import re
import openai
import telebot
import time
import docx
from docx.shared import Pt
from docx import Document
import chardet
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
bot = telebot.TeleBot('6786106072:AAGz8bE4mNyJscUDrvARjYiXOCAyX5FukD4')
openai.api_key = ""


@bot.message_handler(func=lambda message: "презентация" in message.text.lower(), content_types=["text"])
def handle_text(message):
    # Ваш код обработки здесь
    bot.send_message(message.chat.id, "Обработка сообщения с ключевым словом 'презентация'")


    #СТРАНИЦА(1)_заголовок_гл-текст+
    prs = Presentation('PRTIA1.pptx')
    #_заголовок_
    first_slide = prs.slides[1]
    text_box = first_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(0, 0, 0)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[1]
    text_box = first_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER



    #СТРАНИЦА(2)_заголовок_гл-текст
    #_заголовок_
    first_slide = prs.slides[2]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[2]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(1.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER



    #СТРАНИЦА(3)_заголовок_гл-текст
    #_заголовок_
    first_slide = prs.slides[3]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[3]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(1.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER



    #СТРАНИЦА(4)_заголовок_гл-текст
    #_заголовок_
    first_slide = prs.slides[4]
    text_box = first_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(0, 0, 0)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[4]
    text_box = first_slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER



    #СТРАНИЦА(5)_заголовок_гл-текст
    #_заголовок_
    first_slide = prs.slides[5]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[5]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(1.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER



    #СТРАНИЦА(6)_заголовок_гл-текст
    #_заголовок_
    first_slide = prs.slides[6]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(0.3), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс"
    text.font.bold = True
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(28), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER
    #_гл-текст_
    first_slide = prs.slides[6]
    text_box = first_slide.shapes.add_textbox(left=Inches(1), top=Inches(1.5), width=Inches(5), height=Inches(5))
    text_frame = text_box.text_frame
    text = text_frame.add_paragraph()
    text.text = "оформление реферат класс, оформление реферат +по физкультуре, оформление реферат развитие, оформление россия реферат, оформление реферат листьев,оформление  систем реферат, оформление реферат титульный,оформление  титульный лист реферата, оформление виды рефератов, оформление культура оформление реферат, оформление скачать реферат ,оформление реферат +на тему класс ,оформление реферат современные ,оформление основной реферат ,реферат образец,реферат жизни доклад доклад класс доклад +на тему доклад про краткий доклад"
    text.font.name, text.font.size, text.font.color.rgb = 'Times New Roman', Pt(14), RGBColor(255, 255, 255)
    text.alignment = PP_ALIGN.CENTER
    
    '''
    for i in range(5):
        slide = prs.slides[i]
        # вставить изображение
        img_path = f'{i+1}f.PNG' 
        left = Inches(6)
        top = Inches(0.2) 
        width = Inches(5)
        height = Inches(7)
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
    '''
    prs.save('presentation-with-image.pptx')
    bot.send_document(message.chat.id, open('presentation-with-image.pptx', 'rb'))
    file_path = "presentation-with-image.pptx"
    try:
        os.remove(file_path)
        print(f"Файл {file_path} успешно удален.")
    except FileNotFoundError:
        print(f"Файл {file_path} не найден.")
    except Exception as e:
        print(f"Произошла ошибка при удалении файла: {e}")
        
bot.polling(none_stop=True, interval=0)
    
    
